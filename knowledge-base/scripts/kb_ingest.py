#!/usr/bin/env python3
"""kb_ingest — Raw → Processed pipeline for AI Knowledge Engine.

Scans raw/**/unsorted/ for new files, converts them to Markdown when possible,
moves originals to assets/<type>/, computes source hashes, runs NLP enrichment,
estimates complexity, and routes complex materials to review/.

This is the **mechanical** layer: no LLM calls, no semantic decisions.
The agent (or super mode) handles the smart stuff via review queues.

Usage:
  python3 scripts/kb_ingest.py                  # process all raw/*/unsorted/
  python3 scripts/kb_ingest.py path/to/file     # process one file
  python3 scripts/kb_ingest.py --dry-run        # plan, don't execute
  python3 scripts/kb_ingest.py --no-nlp         # skip NLP (faster, less metadata)
  python3 scripts/kb_ingest.py --json           # machine-readable output

Exit codes:
  0 — all files processed
  1 — partial: some files failed (listed in stderr/json)
  2 — environment error (missing dependency, no kb.config.yml)
"""
from __future__ import annotations

import argparse
import datetime as _dt
import json
import logging
import re
import shutil
import sys
import traceback
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Iterable

SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPT_DIR))

import kb_common as kbc  # noqa: E402

logger = logging.getLogger("kb_ingest")

NO_VALUE_PLACEHOLDER = "—"
LONG_BOOK_WORD_THRESHOLD = 25_000
LONG_BOOK_PAGE_THRESHOLD = 80


# ---------------------------------------------------------------------------
# Mapping: extension → asset subfolder & processing strategy
# ---------------------------------------------------------------------------

ASSET_TYPE_BY_EXT: dict[str, str] = {
    ".pdf": "documents",
    ".docx": "documents",
    ".doc": "documents",
    ".txt": "documents",
    ".md": "documents",
    ".rtf": "documents",
    ".odt": "documents",
    ".pptx": "presentations",
    ".ppt": "presentations",
    ".odp": "presentations",
    ".xlsx": "documents",
    ".xls": "documents",
    ".csv": "documents",
    ".ods": "documents",
    ".png": "images",
    ".jpg": "images",
    ".jpeg": "images",
    ".webp": "images",
    ".gif": "images",
    ".svg": "images",
    ".tif": "images",
    ".tiff": "images",
    ".bmp": "images",
    ".mp3": "media",
    ".wav": "media",
    ".m4a": "media",
    ".flac": "media",
    ".ogg": "media",
    ".oga": "media",
    ".aac": "media",
    ".wma": "media",
    ".opus": "media",
    ".mp4": "media",
    ".mov": "media",
    ".webm": "media",
    ".mkv": "media",
    ".avi": "media",
    ".m4v": "media",
    ".mpg": "media",
    ".mpeg": "media",
    ".zip": "archives",
    ".tar": "archives",
    ".tar.gz": "archives",
    ".tgz": "archives",
    ".rar": "archives",
}

PROCESSING_STRATEGY_BY_EXT: dict[str, str] = {
    ".md": "passthrough",
    ".txt": "passthrough",
    ".rtf": "rtf",
    ".docx": "docx",
    ".pdf": "pdf",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
    ".csv": "csv",
    # Images → OCR (rapidocr/tesseract; see kb_ocr.py)
    ".png": "ocr",
    ".jpg": "ocr",
    ".jpeg": "ocr",
    ".webp": "ocr",
    ".tif": "ocr",
    ".tiff": "ocr",
    ".bmp": "ocr",
    # Audio/Video → STT (faster-whisper by default; see kb_stt.py)
    ".mp3": "stt",
    ".wav": "stt",
    ".m4a": "stt",
    ".flac": "stt",
    ".ogg": "stt",
    ".oga": "stt",
    ".aac": "stt",
    ".wma": "stt",
    ".opus": "stt",
    ".mp4": "stt",
    ".mov": "stt",
    ".webm": "stt",
    ".mkv": "stt",
    ".avi": "stt",
    ".m4v": "stt",
    ".mpg": "stt",
    ".mpeg": "stt",
    # Archives → unpack into raw/unsorted for re-ingestion (see kb_ingest._run_archive)
    # Note: bare ".gz" is not listed — gzip-only streams need a container (.tar.gz).
    ".zip": "archive",
    ".tar": "archive",
    ".tar.gz": "archive",
    ".tgz": "archive",
    ".rar": "archive",
}


# ---------------------------------------------------------------------------
# Result types
# ---------------------------------------------------------------------------


@dataclass
class IngestResult:
    source: str                 # original path (relative to root)
    asset_path: str | None = None
    processed_path: str | None = None
    metadata_path: str | None = None
    nlp_meta_path: str | None = None
    routed_to: str = ""          # "knowledge" | "review/needs-ai-decision/" | "skipped"
    review_reason: str = ""
    is_surprise: bool = True
    complexity: float = 0.0
    importance: int | None = None
    source_hash: str = ""
    success: bool = True
    error: str = ""
    skipped: bool = False
    skip_reason: str = ""


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _file_ext(path: Path) -> str:
    """Return the processing extension, treating ``.tar.gz`` as one unit."""
    name = path.name.lower()
    if name.endswith(".tar.gz"):
        return ".tar.gz"
    return path.suffix.lower()


def _detect_asset_type(ext: str) -> str:
    return ASSET_TYPE_BY_EXT.get(ext.lower(), "documents")


def _detect_strategy(ext: str) -> str:
    return PROCESSING_STRATEGY_BY_EXT.get(ext.lower(), "unknown")


def _is_under_unsorted(src: Path, root: Path) -> bool:
    """True when *src* lives under ``raw/**/unsorted/`` (ingest intake)."""
    try:
        rel = src.resolve().relative_to(root.resolve())
    except (OSError, ValueError):
        return False
    parts = rel.parts
    return len(parts) >= 3 and parts[0] == "raw" and "unsorted" in parts[:-1]


def _ensure_kb_dirs(root: Path) -> None:
    """Create directory layout if missing. Idempotent."""
    for sub in ("raw", "processed", "knowledge", "assets", "assets-index", "review",
                "interactions", "scripts", "shell"):
        (root / sub).mkdir(exist_ok=True)
    for sub in kbc.RAW_DIRS:
        (root / "raw" / sub / "unsorted").mkdir(parents=True, exist_ok=True)
    for sub in kbc.PROCESSED_DIRS:
        (root / "processed" / sub).mkdir(parents=True, exist_ok=True)
    for sub in kbc.KNOWLEDGE_DIRS:
        (root / "knowledge" / sub).mkdir(parents=True, exist_ok=True)
    for sub in kbc.ASSET_DIRS:
        (root / "assets" / sub).mkdir(parents=True, exist_ok=True)
    for sub in kbc.REVIEW_DIRS:
        (root / "review" / sub).mkdir(parents=True, exist_ok=True)
    for sub in ("sessions",):
        (root / "interactions" / sub).mkdir(parents=True, exist_ok=True)


def _files_in_unsorted(root: Path) -> list[Path]:
    out: list[Path] = []
    raw = root / "raw"
    if not raw.is_dir():
        return out
    for p in raw.rglob("unsorted/*"):
        if p.is_file():
            out.append(p)
    return out


def _existing_source_hashes(root: Path) -> set[str]:
    """Collect already-recorded source_hashes from extracted-metadata/."""
    out: set[str] = set()
    md_dir = root / "processed" / "extracted-metadata"
    if not md_dir.is_dir():
        return out
    for yml in md_dir.glob("*.yml"):
        try:
            import yaml  # type: ignore[import-untyped]

            data = yaml.safe_load(yml.read_text(encoding="utf-8")) or {}
            h = data.get("source_hash")
            if h:
                out.add(h)
        except Exception:  # noqa: BLE001
            continue
    return out


def _route_log_line(result: IngestResult) -> str:
    line = f"Routed: {result.routed_to}"
    if result.review_reason:
        line += f" ({result.review_reason})"
    return line


# ---------------------------------------------------------------------------
# Converters
# ---------------------------------------------------------------------------


def _convert_passthrough(src: Path) -> str:
    """Copy text content as-is."""
    return src.read_text(encoding="utf-8", errors="replace")


def _convert_docx(src: Path) -> str:
    try:
        import docx  # type: ignore[import-untyped]
    except ImportError as e:
        raise RuntimeError("python-docx not installed") from e
    document = docx.Document(str(src))
    parts = []
    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if "heading 1" in style:
            parts.append(f"# {text}")
        elif "heading 2" in style:
            parts.append(f"## {text}")
        elif "heading 3" in style:
            parts.append(f"### {text}")
        else:
            parts.append(text)
    # Tables
    for tbl in document.tables:
        for row in tbl.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            parts.append("| " + " | ".join(cells) + " |")
    return "\n\n".join(parts) + "\n"


def _convert_rtf(src: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore[import-untyped]
    except ImportError as e:
        raise RuntimeError("striprtf not installed") from e
    raw = src.read_text(encoding="utf-8", errors="replace")
    return rtf_to_text(raw).strip() + "\n"


def _convert_pdf(src: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore[import-untyped]
    except ImportError as e:
        raise RuntimeError("pypdf not installed") from e
    reader = PdfReader(str(src))
    parts = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        parts.append(f"## Page {i}\n\n{text.strip()}")
    return "\n\n".join(parts) + "\n"


def _convert_pptx(src: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError as e:
        raise RuntimeError("python-pptx not installed") from e
    prs = Presentation(str(src))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
    return "\n\n".join(parts) + "\n"


def _convert_xlsx(src: Path) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore[import-untyped]
    except ImportError as e:
        raise RuntimeError("openpyxl not installed") from e
    wb = load_workbook(filename=str(src), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(cells):
                parts.append("| " + " | ".join(cells) + " |")
    return "\n\n".join(parts) + "\n"


def _convert_csv(src: Path) -> str:
    import csv

    parts = []
    with open(src, newline="", encoding="utf-8", errors="replace") as fh:
        reader = csv.reader(fh)
        for i, row in enumerate(reader):
            if i == 0:
                parts.append("| " + " | ".join(row) + " |")
                parts.append("|" + "|".join("---" for _ in row) + "|")
            else:
                parts.append("| " + " | ".join(row) + " |")
    return "\n".join(parts) + "\n"


def _convert(strategy: str, src: Path) -> tuple[str | None, str]:
    """Return (markdown_text, processed_subdir).

    markdown_text is None if conversion is unavailable / impractical.
    processed_subdir tells where to drop the result (e.g. "markdown", "ocr").
    """
    try:
        if strategy == "passthrough":
            return _convert_passthrough(src), "markdown"
        if strategy == "rtf":
            return _convert_rtf(src), "markdown"
        if strategy == "docx":
            return _convert_docx(src), "markdown"
        if strategy == "pdf":
            return _convert_pdf(src), "markdown"
        if strategy == "pptx":
            return _convert_pptx(src), "markdown"
        if strategy == "xlsx":
            return _convert_xlsx(src), "tables"
        if strategy == "csv":
            return _convert_csv(src), "tables"
    except RuntimeError:
        # Optional dep missing → caller decides what to do
        raise
    except Exception as e:  # noqa: BLE001
        raise RuntimeError(f"conversion failed: {e}") from e

    # archive / unknown — not handled here (stt/ocr go through dedicated helpers)
    return None, "markdown"


# ---------------------------------------------------------------------------
# Media converters (STT / OCR / archives) — optional, gracefully degrading
# ---------------------------------------------------------------------------


def _run_stt(src: Path, cfg: kbc.KbConfig) -> tuple[str, str, dict]:
    """Transcribe audio/video. Returns (markdown, "transcripts", meta).

    Raises RuntimeError (caught by caller → review) when STT is disabled or no
    backend is available; the message carries an OS-aware install hint.
    """
    import kb_stt

    if not kb_stt.stt_enabled(cfg):
        raise RuntimeError(
            "STT disabled in kb.config.yml (media.stt.enabled: false). "
            "Transcribe manually or enable it."
        )
    try:
        result = kb_stt.transcribe(src, cfg=cfg)
    except kb_stt.SttUnavailable as e:
        raise RuntimeError(str(e)) from e
    return result.markdown, "transcripts", kb_stt.transcript_metadata(result)


def _run_ocr(src: Path, cfg: kbc.KbConfig) -> tuple[str, str]:
    """OCR an image. Returns (markdown, "ocr"). Raises RuntimeError → review."""
    import kb_ocr

    if not kb_ocr.ocr_enabled(cfg):
        raise RuntimeError(
            "OCR disabled in kb.config.yml (media.ocr.enabled: false). "
            "Describe the image manually or enable it."
        )
    try:
        result = kb_ocr.ocr_image(src, cfg=cfg)
    except kb_ocr.OcrUnavailable as e:
        raise RuntimeError(str(e)) from e
    return result.markdown, "ocr"


def _safe_member_name(stem: str, member_path: str) -> str:
    """Flatten an archive member path into a safe, collision-resistant name."""
    cleaned = member_path.replace("\\", "/").strip("/")
    parts = [p for p in cleaned.split("/") if p not in ("", ".", "..")]
    flat = "__".join(parts) if parts else "file"
    return f"{stem}__{flat}"


def _run_archive(src: Path, root: Path, cfg: kbc.KbConfig) -> tuple[str, str, int]:
    """Unpack an archive into raw/unsorted/unsorted/ for re-ingestion.

    Returns (listing_markdown, "markdown", member_count). Raises RuntimeError
    (→ review) for unsupported formats (e.g. .rar without a helper).
    """
    import tarfile
    import zipfile

    if not cfg.archives.get("enabled", True):
        raise RuntimeError("archive unpacking disabled (media.archives.enabled: false)")

    max_files = int(cfg.archives.get("max_files", 200))
    dest = root / "raw" / "unsorted" / "unsorted"
    dest.mkdir(parents=True, exist_ok=True)
    stem = src.stem
    extracted: list[str] = []

    suffix = _file_ext(src)
    if zipfile.is_zipfile(src):
        with zipfile.ZipFile(src) as zf:
            for info in zf.infolist():
                if info.is_dir():
                    continue
                if len(extracted) >= max_files:
                    break
                target = dest / _safe_member_name(stem, info.filename)
                with zf.open(info) as member, target.open("wb") as out:
                    shutil.copyfileobj(member, out)
                extracted.append(info.filename)
    elif tarfile.is_tarfile(src):
        with tarfile.open(src) as tf:
            for member in tf.getmembers():
                if not member.isfile():
                    continue
                if len(extracted) >= max_files:
                    break
                fobj = tf.extractfile(member)
                if fobj is None:
                    continue
                target = dest / _safe_member_name(stem, member.name)
                with fobj, target.open("wb") as out:
                    shutil.copyfileobj(fobj, out)
                extracted.append(member.name)
    else:
        raise RuntimeError(
            f"unsupported archive format '{suffix}'. Extract it manually into "
            "raw/<category>/unsorted/ and re-run ingest "
            "(.rar needs an external tool such as 'unar'; "
            "bare .gz streams are not supported — use .tar.gz / .zip)."
        )

    listing = [f"# Archive: {src.name}", "", f"Unpacked {len(extracted)} file(s) "
               "into `raw/unsorted/unsorted/` for re-ingestion.", ""]
    for name in extracted[:100]:
        listing.append(f"- {name}")
    if len(extracted) > 100:
        listing.append(f"- … and {len(extracted) - 100} more")
    return "\n".join(listing) + "\n", "markdown", len(extracted)


# ---------------------------------------------------------------------------
# Complexity heuristic
# ---------------------------------------------------------------------------


def estimate_complexity(text: str, *, entity_count: int = 0) -> float:
    """Return 0.0–1.0 complexity score (deterministic heuristic)."""
    score = 0.0
    word_count = len(text.split())
    if word_count > 2000:
        score += 0.2
    if word_count > 5000:
        score += 0.1
    if entity_count > 15:
        score += 0.2
    # Heuristic markers of nuance
    lower = text.lower()
    if any(m in lower for m in ("however", "but ", "однако", "впрочем")):
        score += 0.1
    if "|" in text and text.count("|") > 10:  # tables
        score += 0.1
    if any(ch.isdigit() for ch in text):
        score += 0.05
    return min(score, 1.0)


# ---------------------------------------------------------------------------
# NLP enrichment (best-effort; falls back to empty meta if libs missing)
# ---------------------------------------------------------------------------


def nlp_enrich(text: str, cfg: kbc.KbConfig, knowledge_dir: Path) -> dict:
    """Best-effort NLP enrichment. Returns metadata dict.

    If spaCy/RAKE/KeyBERT are missing, returns a minimal stub with empty lists
    so downstream code can still proceed.
    """
    out: dict = {
        "language": cfg.primary_language,
        "entities": [],
        "keywords": [],
        "unresolved_entities": [],
        "suggested_targets": [],
        "complexity": 0.0,
    }

    # ---------- entities (spaCy) ----------
    try:
        import spacy  # type: ignore[import-untyped]
    except ImportError:
        spacy = None  # type: ignore[assignment]

    entities: list[dict] = []
    if spacy is not None:
        try:
            try:
                nlp = spacy.load(cfg.spacy_model)
            except Exception:
                nlp = spacy.load(cfg.spacy_model_fallback)
            doc = nlp(text[:1_000_000])  # cap for safety
            for ent in doc.ents:
                entities.append(
                    {
                        "text": ent.text,
                        "type": ent.label_,
                        "canonical": ent.text,
                        "existing_page": None,
                    }
                )
        except Exception as e:  # noqa: BLE001
            logger.debug("spaCy NER skipped: %s", e)

    # ---------- entity resolution (fuzzy) ----------
    if entities:
        try:
            from difflib import SequenceMatcher

            slugs = kbc.scan_knowledge_slugs(knowledge_dir)
            slug_keys = list(slugs.keys())
            for ent in entities:
                key = kbc.slugify(ent["text"])
                if key in slugs:
                    ent["canonical"] = key
                    ent["existing_page"] = kbc.posix_relpath(
                        slugs[key][0],
                        knowledge_dir.parent,
                    )
                    continue
                best = (None, 0.0)
                for k in slug_keys:
                    ratio = SequenceMatcher(None, key, k).ratio()
                    if ratio > best[1]:
                        best = (k, ratio)
                if best[0] and best[1] >= cfg.fuzzy_match_threshold:
                    ent["canonical"] = best[0]
                    ent["existing_page"] = kbc.posix_relpath(
                        slugs[best[0]][0],
                        knowledge_dir.parent,
                    )
        except Exception as e:  # noqa: BLE001
            logger.debug("entity resolution skipped: %s", e)

    # Deduplicate entities
    seen: set[str] = set()
    deduped: list[dict] = []
    for ent in entities:
        key = (ent["canonical"], ent.get("type", ""))
        if key in seen:
            continue
        seen.add(key)
        deduped.append(ent)
    out["entities"] = deduped
    out["unresolved_entities"] = [
        e["text"] for e in deduped if not e.get("existing_page")
    ]

    # ---------- keywords (KeyBERT preferred, RAKE fallback) ----------
    try:
        from keybert import KeyBERT  # type: ignore[import-untyped]

        kw_model = KeyBERT()
        keybert_kw = kw_model.extract_keywords(
            text, keyphrase_ngram_range=(1, 3), top_n=cfg.keyword_top_n
        )
        out["keywords"] = [
            {"phrase": phrase, "weight": float(weight), "source": "keybert"}
            for phrase, weight in keybert_kw
        ]
    except ImportError:
        try:
            from rake_nltk import Rake  # type: ignore[import-untyped]

            lang = "russian" if cfg.primary_language == "ru" else "english"
            try:
                rake = Rake(language=lang)
            except Exception:
                rake = Rake()
            rake.extract_keywords_from_text(text)
            top = rake.get_ranked_phrases_with_scores()[: cfg.keyword_top_n]
            out["keywords"] = [
                {"phrase": phrase, "weight": float(score), "source": "rake"}
                for score, phrase in top
            ]
        except Exception as e:  # noqa: BLE001
            logger.debug("RAKE keyword extraction skipped: %s", e)
    except Exception as e:  # noqa: BLE001
        logger.debug("KeyBERT keyword extraction skipped: %s", e)

    out["complexity"] = estimate_complexity(text, entity_count=len(deduped))
    return out


# ---------------------------------------------------------------------------
# Pipeline core
# ---------------------------------------------------------------------------


def _select_date_for_filename(src: Path) -> _dt.date:
    """Use the file's modification date as a sane default."""
    try:
        return _dt.datetime.fromtimestamp(src.stat().st_mtime).date()
    except OSError:
        return _dt.date.today()


def process_one(
    src: Path,
    *,
    root: Path,
    cfg: kbc.KbConfig,
    nlp_enabled: bool = True,
    dry_run: bool = False,
    seen_hashes: set[str] | None = None,
) -> IngestResult:
    rel = src.relative_to(root) if src.is_absolute() and src.is_relative_to(root) else src
    result = IngestResult(source=str(rel))

    try:
        if src.resolve().is_relative_to((root / "assets").resolve()):
            return reprocess_asset(
                src,
                root=root,
                cfg=cfg,
                nlp_enabled=nlp_enabled,
                dry_run=dry_run,
            )
    except OSError:
        pass

    if not _is_under_unsorted(src, root):
        result.success = False
        result.error = (
            f"refusing to ingest path outside raw/**/unsorted/: {src}. "
            "Drop the file under raw/<category>/unsorted/ (or omit explicit "
            "paths to scan unsorted automatically)."
        )
        return result

    ext = _file_ext(src)
    if ext in cfg.skip_extensions:
        result.skipped = True
        result.skip_reason = f"extension {ext} in skip_extensions"
        return result

    asset_type = _detect_asset_type(ext)
    strategy = _detect_strategy(ext)

    # Stable filename
    date = _select_date_for_filename(src)
    new_name = kbc.stable_filename(original_name=src.name, date=date)
    asset_target = root / "assets" / asset_type / new_name

    # Hashing
    try:
        source_hash = kbc.compute_source_hash(src)
    except Exception as e:  # noqa: BLE001
        result.success = False
        result.error = f"hashing failed: {e}"
        return result
    result.source_hash = source_hash

    # Idempotency: already processed?
    if seen_hashes is not None and source_hash in seen_hashes:
        result.skipped = True
        result.skip_reason = f"already ingested (hash {source_hash})"
        return result

    if dry_run:
        result.asset_path = str(asset_target.relative_to(root))
        result.routed_to = "dry-run"
        return result

    # Move original to assets/
    asset_target.parent.mkdir(parents=True, exist_ok=True)
    if asset_target.exists():
        # Same hash already there → skip
        if kbc.compute_source_hash(asset_target) == source_hash:
            result.asset_path = str(asset_target.relative_to(root))
            result.skipped = True
            result.skip_reason = "asset already exists with same hash"
            return result
        # Different content under same name → suffix with hash
        suffix = source_hash.split(":")[-1][:8]
        asset_target = asset_target.with_name(
            f"{asset_target.stem}__{suffix}{asset_target.suffix}"
        )
    shutil.move(str(src), str(asset_target))
    result.asset_path = str(asset_target.relative_to(root))

    # Convert
    md_text: str | None = None
    processed_subdir = "markdown"
    transcript_meta: dict = {}
    archive_count: int | None = None
    try:
        if strategy == "stt":
            md_text, processed_subdir, transcript_meta = _run_stt(asset_target, cfg)
        elif strategy == "ocr":
            md_text, processed_subdir = _run_ocr(asset_target, cfg)
        elif strategy == "archive":
            md_text, processed_subdir, archive_count = _run_archive(
                asset_target, root, cfg
            )
        else:
            md_text, processed_subdir = _convert(strategy, asset_target)
    except RuntimeError as e:
        # Optional dep missing, backend unavailable, or conversion failed
        result.review_reason = f"conversion unavailable: {e}"

    processed_path: Path | None = None
    if md_text is not None:
        processed_path = (
            root / "processed" / processed_subdir / f"{asset_target.stem}.md"
        )
        processed_path.parent.mkdir(parents=True, exist_ok=True)
        processed_path.write_text(md_text, encoding="utf-8")
        result.processed_path = str(processed_path.relative_to(root))

    # NLP
    nlp_meta: dict = {}
    if nlp_enabled and cfg.nlp_enabled and md_text:
        try:
            nlp_meta = nlp_enrich(md_text, cfg, knowledge_dir=root / "knowledge")
            nlp_path = (
                root / "processed" / "nlp-meta" / f"{asset_target.stem}.yml"
            )
            nlp_path.parent.mkdir(parents=True, exist_ok=True)
            import yaml  # type: ignore[import-untyped]

            nlp_path.write_text(
                yaml.safe_dump(nlp_meta, allow_unicode=True, sort_keys=False),
                encoding="utf-8",
            )
            result.nlp_meta_path = str(nlp_path.relative_to(root))
        except Exception as e:  # noqa: BLE001
            logger.debug("NLP enrichment skipped for %s: %s", asset_target.name, e)

    complexity = nlp_meta.get("complexity") if nlp_meta else (
        estimate_complexity(md_text or "")
    )
    result.complexity = float(complexity or 0.0)

    # Detect long reference books (PDFs/EPUBs/DOCX with >25k words)
    long_book = _looks_like_long_book(ext=ext, processed_text=md_text)

    # Routing
    if md_text is None:
        # Cannot auto-process → review
        review_dir = root / "review" / "needs-ai-decision"
        review_pkg = review_dir / f"{asset_target.stem}.md"
        review_pkg.parent.mkdir(parents=True, exist_ok=True)
        review_pkg.write_text(
            _build_review_package(
                asset_path=result.asset_path,
                processed_path=None,
                nlp_meta=nlp_meta,
                reason=result.review_reason or "no automatic conversion available",
                long_book_hint=long_book,
            ),
            encoding="utf-8",
        )
        result.routed_to = "review/needs-ai-decision/"
    elif result.complexity >= cfg.complexity_threshold or long_book:
        review_dir = root / "review" / "needs-ai-decision"
        review_pkg = review_dir / f"{asset_target.stem}.md"
        review_pkg.parent.mkdir(parents=True, exist_ok=True)
        if long_book:
            reason = "looks like a long-form reference book (>=25k words)"
            result.review_reason = reason
        else:
            reason = (
                f"complexity {result.complexity:.2f} >= threshold "
                f"{cfg.complexity_threshold}"
            )
            result.review_reason = (
                f"complexity {result.complexity:.2f} >= {cfg.complexity_threshold}"
            )
        review_pkg.write_text(
            _build_review_package(
                asset_path=result.asset_path,
                processed_path=result.processed_path,
                nlp_meta=nlp_meta,
                reason=reason,
                long_book_hint=long_book,
            ),
            encoding="utf-8",
        )
        result.routed_to = "review/needs-ai-decision/"
    else:
        result.routed_to = "processed/"

    # Metadata
    metadata_path = (
        root / "processed" / "extracted-metadata" / f"{asset_target.stem}.yml"
    )
    metadata_path.parent.mkdir(parents=True, exist_ok=True)
    today = _dt.date.today().isoformat()
    metadata = {
        "original_filename": src.name,
        "stable_filename": asset_target.name,
        "asset_path": result.asset_path,
        "processed_path": result.processed_path,
        "source_hash": source_hash,
        "file_type": ext.lstrip("."),
        "asset_type": asset_type,
        "strategy": strategy,
        "processing_date": kbc.now_iso(),
        "extracted_at": today,
        "valid_from": today,
        "lifecycle": "evolving",
        "confidence": "medium",
        "complexity": result.complexity,
        "is_surprise": True,
        "surprise_engine": "python",
        "long_book_hint": long_book,
        "needs_ai_review": result.routed_to.startswith("review/"),
        "review_reason": result.review_reason,
        "nlp_meta_path": result.nlp_meta_path,
    }
    if transcript_meta:
        metadata.update(transcript_meta)
    if archive_count is not None:
        metadata["archive_extracted_files"] = archive_count
    import yaml  # type: ignore[import-untyped]

    metadata_path.write_text(
        yaml.safe_dump(metadata, allow_unicode=True, sort_keys=False),
        encoding="utf-8",
    )
    result.metadata_path = str(metadata_path.relative_to(root))

    # Update assets-index/<type>.md
    _upsert_assets_index(root, asset_type, asset_target, processed_path)

    # Log
    kbc.append_log(
        operation="ingest",
        title=f"{asset_target.name}",
        details=[
            f"Source hash: {source_hash}",
            f"Asset: {result.asset_path}",
            f"Processed: {result.processed_path or NO_VALUE_PLACEHOLDER}",
            f"NLP meta: {result.nlp_meta_path or NO_VALUE_PLACEHOLDER}",
            f"Complexity: {result.complexity:.2f}",
            *(
                [f"Transcribed: {transcript_meta.get('stt_backend')} "
                 f"({transcript_meta.get('stt_language')}, "
                 f"{transcript_meta.get('stt_segments')} segments)"]
                if transcript_meta else []
            ),
            *(
                [f"Archive: unpacked {archive_count} file(s) into raw/unsorted/"]
                if archive_count is not None else []
            ),
            _route_log_line(result),
        ],
        root=root,
    )
    return result


def _build_review_package(
    *,
    asset_path: str | None,
    processed_path: str | None,
    nlp_meta: dict,
    reason: str,
    long_book_hint: bool = False,
) -> str:
    lines = [f"# AI Review: {Path(asset_path or 'unknown').name}", ""]
    lines.append("## Source")
    lines.append("")
    lines.append(f"- Asset: `{asset_path}`")
    if processed_path:
        lines.append(f"- Processed: `{processed_path}`")
    lines.append(f"- Reason: {reason}")
    if long_book_hint:
        lines.append("")
        lines.append("## ⚠️ Likely long-form reference book")
        lines.append("")
        lines.append(
            "This file looks like a published book or long manual "
            "(many words, many pages, full prose). Before extracting prose into "
            "`knowledge/`, see **03_PIPELINE.md → Handling long reference materials**."
        )
        lines.append("")
        lines.append("Recommended flow:")
        lines.append(
            "1. Confirm the original is in `assets/documents/` "
            "(it should be — pipeline already moved it there)"
        )
        lines.append(
            "2. Add a one-line entry in `assets-index/documents.md` describing the book"
        )
        lines.append(
            "3. Ask the user: \"Want me to draft a 'takeaways for you' note "
            "based on this book? I'll keep it in your own words.\""
        )
        lines.append(
            "4. If yes — write 5-15 bullets in "
            "`knowledge/principles/<book-slug>-takeaways.md` with `source:` frontmatter "
            "pointing back to the asset"
        )
        lines.append(
            "5. Update or create `knowledge/principles/<role>-bookshelf.md` "
            "with the entry"
        )
        lines.append(
            "6. **Do NOT** copy the book's prose into `knowledge/voice/` — it pollutes "
            "the user's own voice. Capture *rules* and *takeaways*, not text."
        )
        lines.append("")
    if nlp_meta:
        lines.append("## NLP enrichment")
        ents = nlp_meta.get("entities", [])
        if ents:
            lines.append(f"- Entities: {len(ents)}")
            for e in ents[:10]:
                resolved = e.get("existing_page") or "—"
                lines.append(f"  - `{e['text']}` ({e.get('type', '?')}) → {resolved}")
        kws = nlp_meta.get("keywords", [])
        if kws:
            lines.append(f"- Keywords (top): {', '.join(k['phrase'] for k in kws[:10])}")
    lines.append("")
    lines.append("## What to do")
    lines.append("")
    if long_book_hint:
        lines.append(
            "**Use the long-book flow above.** The standard steps below apply "
            "only if you've decided not to treat this as a reference book."
        )
        lines.append("")
    lines.append("1. Read the processed Markdown (if any) and the original asset.")
    lines.append("2. Extract durable knowledge into `knowledge/<category>/<slug>.md`")
    lines.append("3. Add frontmatter: `source`, `extracted_at`, `tags`, `lifecycle`, `importance`.")
    lines.append("4. Wikilink to related pages with `[[slug]]`.")
    lines.append("5. Delete this review package once processed.")
    lines.append("")
    return "\n".join(lines)


def _looks_like_long_book(*, ext: str, processed_text: str | None) -> bool:
    """Heuristic: does this look like a copyrighted long-form book / manual?

    Triggers when *all* of these hold:
    - File extension is a typical long-form document format
    - Extracted text is either very long by words, or a PDF with many page
      markers even if the extractor collapsed intra-line spacing

    Returns False for the user's own short drafts, articles, and notes.
    """
    normalized_ext = ext.lower()
    if normalized_ext not in {".pdf", ".epub", ".docx", ".doc", ".odt"}:
        return False
    if not processed_text:
        return False

    word_count = len(processed_text.split())
    if word_count >= LONG_BOOK_WORD_THRESHOLD:
        return True

    if normalized_ext == ".pdf":
        page_markers = len(re.findall(r"(?m)^## Page \d+\b", processed_text))
        if page_markers >= LONG_BOOK_PAGE_THRESHOLD:
            return True

    return False


def _upsert_assets_index(
    root: Path, asset_type: str, asset_path: Path, processed_path: Path | None
) -> None:
    """Write or replace a single asset block in assets-index/<type>.md."""
    index_file = root / "assets-index" / f"{asset_type}.md"
    index_file.parent.mkdir(parents=True, exist_ok=True)
    if not index_file.exists():
        index_file.write_text(
            f"# {asset_type.capitalize()}\n\n"
            "Asset descriptions for files in `assets/" + asset_type + "/`.\n",
            encoding="utf-8",
        )

    rel_asset = asset_path.relative_to(root).as_posix()
    rel_processed = (
        processed_path.relative_to(root).as_posix() if processed_path else None
    )
    block_lines = [
        f"## {asset_path.stem}",
        "",
        f"- Type: {asset_type}",
        f"- Original: `{rel_asset}`",
        f"- Converted: `{rel_processed}`" if rel_processed else "- Converted: none yet",
        "- Description: (to be filled by AI agent during review)",
        "",
    ]
    block_text = "\n".join(block_lines) + "\n"
    heading = f"## {asset_path.stem}"
    text = index_file.read_text(encoding="utf-8")
    pattern = rf"\n{re.escape(heading)}\n(?:.*?)(?=\n## |\Z)"

    if heading in text:
        updated = re.sub(pattern, lambda _: "\n" + block_text, text, flags=re.S)
        index_file.write_text(updated, encoding="utf-8")
        return

    with index_file.open("a", encoding="utf-8") as fh:
        if text and not text.endswith("\n"):
            fh.write("\n")
        fh.write("\n" + block_text)


def reprocess_asset(
    src: Path,
    *,
    root: Path,
    cfg: kbc.KbConfig,
    nlp_enabled: bool = True,
    dry_run: bool = False,
) -> IngestResult:
    """Re-run conversion, metadata, and review routing for an existing asset."""
    rel = src.relative_to(root) if src.is_absolute() and src.is_relative_to(root) else src
    result = IngestResult(source=str(rel), asset_path=str(rel))

    ext = _file_ext(src)
    asset_type = _detect_asset_type(ext)
    strategy = _detect_strategy(ext)
    stem = src.stem

    try:
        source_hash = kbc.compute_source_hash(src)
    except Exception as e:  # noqa: BLE001
        result.success = False
        result.error = f"hashing failed: {e}"
        return result
    result.source_hash = source_hash

    if dry_run:
        result.routed_to = "dry-run"
        return result

    md_text: str | None = None
    processed_subdir = "markdown"
    transcript_meta: dict = {}
    archive_count: int | None = None
    try:
        if strategy == "stt":
            md_text, processed_subdir, transcript_meta = _run_stt(src, cfg)
        elif strategy == "ocr":
            md_text, processed_subdir = _run_ocr(src, cfg)
        elif strategy == "archive":
            md_text, processed_subdir, archive_count = _run_archive(src, root, cfg)
        else:
            md_text, processed_subdir = _convert(strategy, src)
    except RuntimeError as e:
        result.review_reason = f"conversion unavailable: {e}"

    processed_path: Path | None = None
    if md_text is not None:
        processed_path = root / "processed" / processed_subdir / f"{stem}.md"
        processed_path.parent.mkdir(parents=True, exist_ok=True)
        processed_path.write_text(md_text, encoding="utf-8")
        result.processed_path = str(processed_path.relative_to(root))

    nlp_meta: dict = {}
    if nlp_enabled and cfg.nlp_enabled and md_text:
        try:
            nlp_meta = nlp_enrich(md_text, cfg, knowledge_dir=root / "knowledge")
            nlp_path = root / "processed" / "nlp-meta" / f"{stem}.yml"
            nlp_path.parent.mkdir(parents=True, exist_ok=True)
            import yaml  # type: ignore[import-untyped]

            nlp_path.write_text(
                yaml.safe_dump(nlp_meta, allow_unicode=True, sort_keys=False),
                encoding="utf-8",
            )
            result.nlp_meta_path = str(nlp_path.relative_to(root))
        except Exception as e:  # noqa: BLE001
            logger.debug("NLP enrichment skipped for %s: %s", src.name, e)

    complexity = nlp_meta.get("complexity") if nlp_meta else estimate_complexity(md_text or "")
    result.complexity = float(complexity or 0.0)
    long_book = _looks_like_long_book(ext=ext, processed_text=md_text)

    review_pkg = root / "review" / "needs-ai-decision" / f"{stem}.md"
    if md_text is None:
        review_pkg.parent.mkdir(parents=True, exist_ok=True)
        review_pkg.write_text(
            _build_review_package(
                asset_path=result.asset_path,
                processed_path=None,
                nlp_meta=nlp_meta,
                reason=result.review_reason or "no automatic conversion available",
                long_book_hint=long_book,
            ),
            encoding="utf-8",
        )
        result.routed_to = "review/needs-ai-decision/"
    elif result.complexity >= cfg.complexity_threshold or long_book:
        review_pkg.parent.mkdir(parents=True, exist_ok=True)
        if long_book:
            reason = "looks like a long-form reference book (>=25k words)"
            result.review_reason = reason
        else:
            reason = (
                f"complexity {result.complexity:.2f} >= threshold "
                f"{cfg.complexity_threshold}"
            )
            result.review_reason = (
                f"complexity {result.complexity:.2f} >= {cfg.complexity_threshold}"
            )
        review_pkg.write_text(
            _build_review_package(
                asset_path=result.asset_path,
                processed_path=result.processed_path,
                nlp_meta=nlp_meta,
                reason=reason,
                long_book_hint=long_book,
            ),
            encoding="utf-8",
        )
        result.routed_to = "review/needs-ai-decision/"
    else:
        if review_pkg.exists():
            review_pkg.unlink()
        result.routed_to = "processed/"
        result.review_reason = ""

    metadata_path = root / "processed" / "extracted-metadata" / f"{stem}.yml"
    metadata_path.parent.mkdir(parents=True, exist_ok=True)
    today = _dt.date.today().isoformat()

    original_filename = src.name
    if metadata_path.exists():
        try:
            import yaml  # type: ignore[import-untyped]

            previous = yaml.safe_load(metadata_path.read_text(encoding="utf-8")) or {}
            original_filename = previous.get("original_filename", original_filename)
        except Exception:  # noqa: BLE001
            pass

    metadata = {
        "original_filename": original_filename,
        "stable_filename": src.name,
        "asset_path": result.asset_path,
        "processed_path": result.processed_path,
        "source_hash": source_hash,
        "file_type": ext.lstrip("."),
        "asset_type": asset_type,
        "strategy": strategy,
        "processing_date": kbc.now_iso(),
        "extracted_at": today,
        "valid_from": today,
        "lifecycle": "evolving",
        "confidence": "medium",
        "complexity": result.complexity,
        "is_surprise": True,
        "surprise_engine": "python",
        "long_book_hint": long_book,
        "needs_ai_review": result.routed_to.startswith("review/"),
        "review_reason": result.review_reason,
        "nlp_meta_path": result.nlp_meta_path,
    }
    if transcript_meta:
        metadata.update(transcript_meta)
    if archive_count is not None:
        metadata["archive_extracted_files"] = archive_count
    import yaml  # type: ignore[import-untyped]

    metadata_path.write_text(
        yaml.safe_dump(metadata, allow_unicode=True, sort_keys=False),
        encoding="utf-8",
    )
    result.metadata_path = str(metadata_path.relative_to(root))

    _upsert_assets_index(root, asset_type, src, processed_path)

    kbc.append_log(
        operation="reprocess",
        title=f"{src.name}",
        details=[
            f"Source hash: {source_hash}",
            f"Asset: {result.asset_path}",
            f"Processed: {result.processed_path or NO_VALUE_PLACEHOLDER}",
            f"NLP meta: {result.nlp_meta_path or NO_VALUE_PLACEHOLDER}",
            f"Complexity: {result.complexity:.2f}",
            *(
                [f"Transcribed: {transcript_meta.get('stt_backend')} "
                 f"({transcript_meta.get('stt_language')}, "
                 f"{transcript_meta.get('stt_segments')} segments)"]
                if transcript_meta else []
            ),
            *(
                [f"Archive: unpacked {archive_count} file(s) into raw/unsorted/"]
                if archive_count is not None else []
            ),
            _route_log_line(result),
        ],
        root=root,
    )
    return result


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="AI Knowledge Engine — ingest pipeline")
    parser.add_argument("paths", nargs="*", type=Path, help="Specific files to process")
    parser.add_argument("--root", type=Path, default=None)
    parser.add_argument("--dry-run", action="store_true")
    parser.add_argument("--no-nlp", action="store_true")
    parser.add_argument("--json", action="store_true")
    parser.add_argument(
        "--verbose", action="store_true", help="Show debug logs (skipped NLP, etc.)"
    )
    parser.add_argument(
        "--init-dirs",
        action="store_true",
        help="Just create directory structure and exit",
    )
    args = parser.parse_args(argv)

    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.WARNING,
        format="[%(name)s] %(levelname)s: %(message)s",
    )

    root = args.root or kbc.find_kb_root()
    cfg = kbc.load_config(root)

    _ensure_kb_dirs(root)
    if args.init_dirs:
        print(f"Directory structure ensured under {root}")
        return 0

    if args.paths:
        files = [p for p in args.paths if p.is_file()]
    else:
        files = _files_in_unsorted(root)

    if not files:
        if args.json:
            print(json.dumps({"processed": [], "errors": [], "skipped": []}))
        else:
            print("No files in raw/*/unsorted/. Drop materials there and re-run.")
        return 0

    seen_hashes = _existing_source_hashes(root)
    results: list[IngestResult] = []
    for f in files:
        try:
            r = process_one(
                f,
                root=root,
                cfg=cfg,
                nlp_enabled=not args.no_nlp,
                dry_run=args.dry_run,
                seen_hashes=seen_hashes,
            )
        except Exception as e:  # noqa: BLE001
            r = IngestResult(
                source=str(f),
                success=False,
                error=f"{e}\n{traceback.format_exc(limit=2)}",
            )
        results.append(r)
        if r.source_hash:
            seen_hashes.add(r.source_hash)

    failures = [r for r in results if not r.success]
    if args.json:
        print(
            json.dumps(
                {
                    "processed": [asdict(r) for r in results if r.success and not r.skipped],
                    "skipped": [asdict(r) for r in results if r.skipped],
                    "errors": [asdict(r) for r in failures],
                },
                ensure_ascii=False,
                indent=2,
            )
        )
    else:
        for r in results:
            if r.skipped:
                print(f"⏭  {r.source} ({r.skip_reason})")
            elif not r.success:
                print(f"❌ {r.source}: {r.error}", file=sys.stderr)
            else:
                badge = "📥" if r.routed_to.startswith("review/") else "✅"
                print(
                    f"{badge} {r.source} → {r.routed_to} "
                    f"(complexity {r.complexity:.2f})"
                )

    return 1 if failures else 0


if __name__ == "__main__":
    raise SystemExit(main())
